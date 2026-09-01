import base64
import binascii
import io
from typing import Any, Dict, List, Optional, Tuple
import docx
import httpx
import pandas as pd
import pptx

MAX_IMAGE_BYTES = 25 * 1024 * 1024  # 25MB
MAX_ATTACHMENTS = 5
FETCH_TIMEOUT_SECONDS = 8.0

# Supported native multimodal types for Gemini Part.from_bytes
NATIVE_MULTIMODAL_MIMES = {
    "image/jpeg",
    "image/png",
    "image/webp",
    "image/heic",
    "application/pdf",
    "video/mp4",
    "video/quicktime",
    "video/webm",
    "video/mpeg",
}


def extract_text_from_doc(
    file_bytes: bytes, mime_type: str, filename: str = "doc"
) -> str:
    """Extracts text or markdown table representations from non-native files."""
    try:
        # Plain text / CSV / JSON / XML
        if mime_type in {
            "text/plain",
            "application/json",
            "application/xml",
            "text/xml",
        }:
            return file_bytes.decode("utf-8", errors="ignore")

        if mime_type == "text/csv" or filename.endswith(".csv"):
            df = pd.read_csv(io.BytesIO(file_bytes))
            return df.to_markdown(index=False)

        # Excel Spreadsheets
        if (
            "spreadsheet" in mime_type
            or "excel" in mime_type
            or filename.endswith((".xlsx", ".xls"))
        ):
            excel_file = pd.ExcelFile(io.BytesIO(file_bytes))
            sheet_texts = []
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(excel_file, sheet_name=sheet_name)
                sheet_texts.append(
                    f"### Sheet: {sheet_name}\n" + df.to_markdown(index=False)
                )
            return "\n\n".join(sheet_texts)

        # Word Documents
        if (
            "wordprocessingml" in mime_type
            or "msword" in mime_type
            or filename.endswith((".docx", ".doc"))
        ):
            doc = docx.Document(io.BytesIO(file_bytes))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            table_texts = []
            for table in doc.tables:
                rows = []
                for row in table.rows:
                    rows.append(" | ".join(cell.text.strip() for cell in row.cells))
                table_texts.append("\n".join(rows))

            content_parts = []
            if paragraphs:
                content_parts.append("\n".join(paragraphs))
            if table_texts:
                content_parts.append("Tables:\n" + "\n\n".join(table_texts))
            return "\n\n".join(content_parts)

        # PowerPoint Presentations
        if "presentation" in mime_type or filename.endswith((".pptx", ".ppt")):
            prs = pptx.Presentation(io.BytesIO(file_bytes))
            slide_texts = []
            for idx, slide in enumerate(prs.slides, 1):
                texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text.strip())
                if texts:
                    slide_texts.append(f"### Slide {idx}:\n" + "\n".join(texts))
            return "\n\n".join(slide_texts)

        # Raw Binary
        if mime_type == "application/octet-stream" or filename.endswith(".bin"):
            chunk = file_bytes[:1024]
            hex_data = binascii.hexlify(chunk, b" ", 2).decode("ascii")
            return (
                f"BINARY HEX DUMP (First 1KB of {len(file_bytes)} bytes):\n{hex_data}"
            )

    except Exception as e:
        return f"[Error extracting text from {filename} ({mime_type}): {str(e)}]"

    return ""


def decode_attachment(attachment: Dict[str, Any]) -> Optional[Tuple[bytes, str, str]]:
    """Decodes attachment into (bytes, mime_type, name)."""
    if not isinstance(attachment, dict):
        return None
    mime_type = str(attachment.get("mime_type") or "").lower()
    name = str(attachment.get("name") or "file")
    data = attachment.get("data")
    if not isinstance(data, str) or not data:
        return None

    if "," in data and data.lstrip().startswith("data:"):
        data = data.split(",", 1)[1]

    try:
        decoded = base64.b64decode(data, validate=False)
    except (binascii.Error, ValueError):
        return None

    if not decoded or len(decoded) > MAX_IMAGE_BYTES:
        return None
    return decoded, mime_type, name


def decode_attachments(
    attachments: Optional[List[Dict[str, Any]]],
) -> List[Tuple[bytes, str, str]]:
    """Decodes all attachments, capping at MAX_ATTACHMENTS."""
    if not attachments:
        return []
    decoded = []
    for attachment in attachments[:MAX_ATTACHMENTS]:
        result = decode_attachment(attachment)
        if result:
            decoded.append(result)
    return decoded


async def fetch_photo(url: str) -> Optional[Tuple[bytes, str]]:
    if not url or not url.startswith("https://"):
        return None
    try:
        async with httpx.AsyncClient(timeout=FETCH_TIMEOUT_SECONDS) as client:
            response = await client.get(url)
            if response.status_code != 200 or len(response.content) > MAX_IMAGE_BYTES:
                return None
            mime_type = (
                response.headers.get("content-type", "image/jpeg")
                .split(";")[0]
                .strip()
                .lower()
            )
            return response.content, mime_type
    except Exception:
        return None


async def fetch_photos(urls: List[str]) -> List[Tuple[bytes, str]]:
    images = []
    for url in urls[:MAX_ATTACHMENTS]:
        result = await fetch_photo(url)
        if result:
            images.append(result)
    return images


def wants_body_fat_estimate(message: str) -> bool:
    lowered = (message or "").lower()
    triggers = (
        "body fat",
        "bodyfat",
        "bf%",
        "physique",
        "lean",
        "leaner",
        "how do i look",
        "estimate my",
        "progress photo",
        "shredded",
        "conditioning",
    )
    return any(trigger in lowered for trigger in triggers)
